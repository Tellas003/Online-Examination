from django.shortcuts import render, redirect, reverse, HttpResponseRedirect
from . import forms, models
from django.db.models import Sum
from django.contrib.auth.models import Group, User
from django.contrib.auth.decorators import login_required, user_passes_test
from django.conf import settings
from datetime import date, timedelta
from django.core.mail import send_mail
from teacher import models as TMODEL
from student import models as SMODEL
from teacher import forms as TFORM
from student import forms as SFORM
from exam import models as QMODEL
import pandas as pd
from docx import Document
from pptx import Presentation


# ----------------- Common Views -----------------

def home_view(request):
    if request.user.is_authenticated:
        return HttpResponseRedirect('afterlogin')
    return render(request, 'exam/index.html')


def is_teacher(user):
    return user.groups.filter(name='TEACHER').exists()


def is_student(user):
    return user.groups.filter(name='STUDENT').exists()


def afterlogin_view(request):
    if is_student(request.user):
        return redirect('student/student-dashboard')

    elif is_teacher(request.user):
        accountapproval = TMODEL.Teacher.objects.filter(user_id=request.user.id, status=True)
        if accountapproval:
            return redirect('teacher/teacher-dashboard')
        else:
            return render(request, 'teacher/teacher_wait_for_approval.html')

    else:
        return redirect('admin-dashboard')


# ----------------- Admin Views -----------------

def adminclick_view(request):
    if request.user.is_authenticated:
        return HttpResponseRedirect('afterlogin')
    return HttpResponseRedirect('adminlogin')


@login_required(login_url='adminlogin')
def admin_dashboard_view(request):
    context = {
        'total_student': SMODEL.Student.objects.all().count(),
        'total_teacher': TMODEL.Teacher.objects.filter(status=True).count(),
        'total_course': models.Course.objects.all().count(),
        'total_question': models.Question.objects.all().count(),
    }
    return render(request, 'exam/admin_dashboard.html', context)


@login_required(login_url='adminlogin')
def admin_teacher_view(request):
    context = {
        'total_teacher': TMODEL.Teacher.objects.filter(status=True).count(),
        'pending_teacher': TMODEL.Teacher.objects.filter(status=False).count(),
        'salary': TMODEL.Teacher.objects.filter(status=True).aggregate(Sum('salary'))['salary__sum'],
    }
    return render(request, 'exam/admin_teacher.html', context)


@login_required(login_url='adminlogin')
def admin_view_teacher_view(request):
    teachers = TMODEL.Teacher.objects.filter(status=True)
    return render(request, 'exam/admin_view_teacher.html', {'teachers': teachers})


@login_required(login_url='adminlogin')
def update_teacher_view(request, pk):
    teacher = TMODEL.Teacher.objects.get(id=pk)
    user = User.objects.get(id=teacher.user_id)
    userForm = TFORM.TeacherUserForm(instance=user)
    teacherForm = TFORM.TeacherForm(instance=teacher)
    if request.method == 'POST':
        userForm = TFORM.TeacherUserForm(request.POST, instance=user)
        teacherForm = TFORM.TeacherForm(request.POST, request.FILES, instance=teacher)
        if userForm.is_valid() and teacherForm.is_valid():
            user = userForm.save()
            user.set_password(user.password)
            user.save()
            teacherForm.save()
            return redirect('admin-view-teacher')
    return render(request, 'exam/update_teacher.html', {'userForm': userForm, 'teacherForm': teacherForm})


@login_required(login_url='adminlogin')
def delete_teacher_view(request, pk):
    teacher = TMODEL.Teacher.objects.get(id=pk)
    user = User.objects.get(id=teacher.user_id)
    user.delete()
    teacher.delete()
    return HttpResponseRedirect('/admin-view-teacher')


@login_required(login_url='adminlogin')
def admin_view_pending_teacher_view(request):
    teachers = TMODEL.Teacher.objects.filter(status=False)
    return render(request, 'exam/admin_view_pending_teacher.html', {'teachers': teachers})


@login_required(login_url='adminlogin')
def approve_teacher_view(request, pk):
    teacherSalary = forms.TeacherSalaryForm()
    if request.method == 'POST':
        teacherSalary = forms.TeacherSalaryForm(request.POST)
        if teacherSalary.is_valid():
            teacher = TMODEL.Teacher.objects.get(id=pk)
            teacher.salary = teacherSalary.cleaned_data['salary']
            teacher.status = True
            teacher.save()
        return HttpResponseRedirect('/admin-view-pending-teacher')
    return render(request, 'exam/salary_form.html', {'teacherSalary': teacherSalary})


@login_required(login_url='adminlogin')
def reject_teacher_view(request, pk):
    teacher = TMODEL.Teacher.objects.get(id=pk)
    user = User.objects.get(id=teacher.user_id)
    user.delete()
    teacher.delete()
    return HttpResponseRedirect('/admin-view-pending-teacher')


@login_required(login_url='adminlogin')
def admin_view_teacher_salary_view(request):
    teachers = TMODEL.Teacher.objects.filter(status=True)
    return render(request, 'exam/admin_view_teacher_salary.html', {'teachers': teachers})


@login_required(login_url='adminlogin')
def admin_student_view(request):
    context = {'total_student': SMODEL.Student.objects.all().count()}
    return render(request, 'exam/admin_student.html', context)


@login_required(login_url='adminlogin')
def admin_view_student_view(request):
    students = SMODEL.Student.objects.all()
    return render(request, 'exam/admin_view_student.html', {'students': students})


@login_required(login_url='adminlogin')
def update_student_view(request, pk):
    student = SMODEL.Student.objects.get(id=pk)
    user = User.objects.get(id=student.user_id)
    userForm = SFORM.StudentUserForm(instance=user)
    studentForm = SFORM.StudentForm(instance=student)
    if request.method == 'POST':
        userForm = SFORM.StudentUserForm(request.POST, instance=user)
        studentForm = SFORM.StudentForm(request.POST, request.FILES, instance=student)
        if userForm.is_valid() and studentForm.is_valid():
            user = userForm.save()
            user.set_password(user.password)
            user.save()
            studentForm.save()
            return redirect('admin-view-student')
    return render(request, 'exam/update_student.html', {'userForm': userForm, 'studentForm': studentForm})


@login_required(login_url='adminlogin')
def delete_student_view(request, pk):
    student = SMODEL.Student.objects.get(id=pk)
    user = User.objects.get(id=student.user_id)
    user.delete()
    student.delete()
    return HttpResponseRedirect('/admin-view-student')


# ----------------- Course Views -----------------

@login_required(login_url='adminlogin')
def admin_course_view(request):
    return render(request, 'exam/admin_course.html')


@login_required(login_url='adminlogin')
def admin_add_course_view(request):
    courseForm = forms.CourseForm()
    if request.method == 'POST':
        courseForm = forms.CourseForm(request.POST)
        if courseForm.is_valid():
            courseForm.save()
            return HttpResponseRedirect('/admin-view-course')
    return render(request, 'exam/admin_add_course.html', {'courseForm': courseForm})


@login_required(login_url='adminlogin')
def admin_view_course_view(request):
    courses = models.Course.objects.all()
    return render(request, 'exam/admin_view_course.html', {'courses': courses})


@login_required(login_url='adminlogin')
def delete_course_view(request, pk):
    course = models.Course.objects.get(id=pk)
    course.delete()
    return HttpResponseRedirect('/admin-view-course')


# ----------------- Question Views -----------------

@login_required(login_url='adminlogin')
def admin_question_view(request):
    return render(request, 'exam/admin_question.html')


@login_required(login_url='adminlogin')
def admin_add_question_view(request):
    """Admin can add or upload questions in XLSX, CSV, DOCX, or PPTX formats."""
    questionForm = forms.QuestionForm()
    courses = models.Course.objects.all()

    if request.method == 'POST':
        # ✅ Manual input
        if 'question' in request.POST:
            questionForm = forms.QuestionForm(request.POST, request.FILES)
            if questionForm.is_valid():
                question = questionForm.save(commit=False)
                course_id = request.POST.get('courseID')
                course = models.Course.objects.get(id=course_id)
                question.course = course
                question.marks = 1
                question.save()
                return HttpResponseRedirect('/admin-view-question')

        # ✅ File upload
        elif 'question_file' in request.FILES:
            uploaded_file = request.FILES['question_file']
            ext = uploaded_file.name.split('.')[-1].lower()
            course_id = request.POST.get('courseID')
            course = models.Course.objects.get(id=course_id)

            try:
                # ---------------- XLSX ----------------
                if ext == 'xlsx':
                    df = pd.read_excel(uploaded_file)
                    for _, row in df.iterrows():
                        models.Question.objects.create(
                            course=course,
                            question=row['Question'],
                            option1=row['Option1'],
                            option2=row['Option2'],
                            option3=row['Option3'],
                            option4=row['Option4'],
                            answer=row['Answer'].strip(),
                            marks=1
                        )

                # ---------------- CSV ----------------
                elif ext == 'csv':
                    df = pd.read_csv(uploaded_file)
                    for _, row in df.iterrows():
                        models.Question.objects.create(
                            course=course,
                            question=row['Question'],
                            option1=row['Option1'],
                            option2=row['Option2'],
                            option3=row['Option3'],
                            option4=row['Option4'],
                            answer=row['Answer'].strip(),
                            marks=1
                        )

                # ---------------- DOCX ----------------
                elif ext == 'docx':
                    doc = Document(uploaded_file)
                    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                    for i in range(0, len(lines), 6):
                        try:
                            q = lines[i]
                            o1, o2, o3, o4, ans = lines[i+1:i+6]
                            models.Question.objects.create(
                                course=course,
                                question=q,
                                option1=o1, option2=o2, option3=o3, option4=o4,
                                answer=ans.strip(),
                                marks=1
                            )
                        except Exception:
                            continue

                # ---------------- PPTX ----------------
                elif ext == 'pptx':
                    prs = Presentation(uploaded_file)
                    slides_text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slides_text.append(shape.text.strip())

                    for i in range(0, len(slides_text), 6):
                        try:
                            q = slides_text[i]
                            o1, o2, o3, o4, ans = slides_text[i+1:i+6]
                            models.Question.objects.create(
                                course=course,
                                question=q,
                                option1=o1, option2=o2, option3=o3, option4=o4,
                                answer=ans.strip(),
                                marks=1
                            )
                        except Exception:
                            continue
                

                else:
                    return render(request, 'exam/admin_add_question.html', {
                        'error': 'Unsupported file type. Upload XLSX, CSV, DOCX, or PPTX.',
                        'questionForm': questionForm,
                        'courses': courses
                    })

                return HttpResponseRedirect('/admin-view-question')

            except Exception as e:
                return render(request, 'exam/admin_add_question.html', {
                    'error': f'Error processing file: {e}',
                    'questionForm': questionForm,
                    'courses': courses
                })

    return render(request, 'exam/admin_add_question.html', {
        'questionForm': questionForm,
        'courses': courses
    })


@login_required(login_url='adminlogin')
def admin_view_question_view(request):
    courses = models.Course.objects.all()
    return render(request, 'exam/admin_view_question.html', {'courses': courses})


@login_required(login_url='adminlogin')
def view_question_view(request, pk):
    questions = models.Question.objects.filter(course_id=pk)
    return render(request, 'exam/view_question.html', {'questions': questions})


@login_required(login_url='adminlogin')
def delete_question_view(request, pk):
    question = models.Question.objects.get(id=pk)
    question.delete()
    return HttpResponseRedirect('/admin-view-question')


# ----------------- Marks View -----------------

@login_required(login_url='adminlogin')
def admin_view_student_marks_view(request):
    students = SMODEL.Student.objects.all()
    return render(request, 'exam/admin_view_student_marks.html', {'students': students})


@login_required(login_url='adminlogin')
def admin_view_marks_view(request, pk):
    courses = models.Course.objects.all()
    response = render(request, 'exam/admin_view_marks.html', {'courses': courses})
    response.set_cookie('student_id', str(pk))
    return response


@login_required(login_url='adminlogin')
def admin_check_marks_view(request, pk):
    course = models.Course.objects.get(id=pk)
    student_id = request.COOKIES.get('student_id')
    student = SMODEL.Student.objects.get(id=student_id)
    questions = QMODEL.Question.objects.filter(course=course)
    total_marks = sum(q.marks for q in questions)
    results = models.Result.objects.filter(exam=course, student=student)
    return render(request, 'exam/admin_check_marks.html', {'results': results, 'total_marks': total_marks})


# ----------------- Static Pages -----------------

def aboutus_view(request):
    return render(request, 'exam/aboutus.html')


def contactus_view(request):
    sub = forms.ContactusForm()
    if request.method == 'POST':
        sub = forms.ContactusForm(request.POST)
        if sub.is_valid():
            email = sub.cleaned_data['Email']
            name = sub.cleaned_data['Name']
            message = sub.cleaned_data['Message']
            send_mail(f"{name} || {email}", message, settings.EMAIL_HOST_USER, settings.EMAIL_RECEIVING_USER)
            return render(request, 'exam/contactussuccess.html')
    return render(request, 'exam/contactus.html', {'form': sub})

from django.shortcuts import render
from django.contrib.auth.decorators import login_required
from exam import models as QMODEL
from student import models as SMODEL
from teacher import models as TMODEL



@login_required(login_url='adminlogin')
def admin_teacher_student_marks_view(request):
    teachers = TMODEL.Teacher.objects.filter(status=True)
    selected_teacher_id = request.GET.get('teacher')
    students = []
    student_answers = []

    if selected_teacher_id:
        teacher = TMODEL.Teacher.objects.get(id=selected_teacher_id)
        courses = QMODEL.Course.objects.filter(teacher=teacher)
        selected_course_id = request.GET.get('course')

        if selected_course_id:
            course = QMODEL.Course.objects.get(id=selected_course_id)
            students = SMODEL.Student.objects.all()
            selected_student_id = request.GET.get('student')

            if selected_student_id:
                student = SMODEL.Student.objects.get(id=selected_student_id)
                student_answers = QMODEL.StudentAnswer.objects.filter(
                    question__course=course,
                    student=student
                )

    return render(request, 'exam/admin_teacher_student_marks.html', {
        'teachers': teachers,
        'students': students,
        'student_answers': student_answers,
    })

# ----------------- Admin - Teacher & Student Marks View -----------------

from teacher.models import Teacher
from exam.models import Result, Question, Course

@login_required(login_url='adminlogin')
def admin_teacher_student_marks_view(request):
    from teacher.models import Teacher
    from student.models import Student
    from exam.models import Result

    teachers = Teacher.objects.filter(status=True)

    selected_teacher = request.GET.get('teacher_id')
    students = []
    results = []

    if selected_teacher:
        students = Student.objects.filter(teacher_id=selected_teacher)
        results = Result.objects.filter(student__in=students)

    return render(request, 'exam/admin_teacher_student_marks.html', {
        'teachers': teachers,
        'students': students,
        'results': results,
        'selected_teacher': selected_teacher,
    })
